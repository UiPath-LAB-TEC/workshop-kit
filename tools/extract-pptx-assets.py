#!/usr/bin/env python3
"""Extract screenshots, text and callout geometry from a facilitator deck.

Ran as a per-repo script copy in two workshops before it lived here. Those two
copies diverged, and the older one silently dropped every image that had been
dropped into a PowerPoint content placeholder rather than pasted as a picture --
18 of 33 screenshots on one deck, with no error. That fix is at `is_picture`
below, and it is the reason this file is shared rather than copied.

The two things that were genuinely per-deck in those copies -- which slides
belong to which docs page, and the role keywords -- are now a `--slide-map`
JSON file supplied by the caller. See `schemas/pptx-slide-map.schema.json`.

Needs `python-pptx` and `Pillow`. The `workshop-kit extract-pptx` wrapper
checks for both before invoking this.
"""
from __future__ import annotations

import argparse
import hashlib
import json
import re
from pathlib import Path
from typing import Any

from PIL import Image
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE


EMU_PER_INCH = 914400
UNASSIGNED = "unassigned"


def inches(value: int | None) -> float | None:
    if value is None:
        return None
    return round(value / EMU_PER_INCH, 3)


def shape_rect(shape: Any) -> dict[str, float | None]:
    return {
        "x_in": inches(getattr(shape, "left", None)),
        "y_in": inches(getattr(shape, "top", None)),
        "w_in": inches(getattr(shape, "width", None)),
        "h_in": inches(getattr(shape, "height", None)),
    }


def parse_slide_selector(entries: Any, context: str) -> list[tuple[int, int]]:
    """`[2, "3-18"]` -> `[(2, 2), (3, 18)]`. Ranges are inclusive, as printed."""
    if entries is None:
        return []
    if not isinstance(entries, list):
        raise ValueError(f"{context}: \"slides\" must be a list")

    ranges: list[tuple[int, int]] = []
    for entry in entries:
        if isinstance(entry, bool):
            raise ValueError(f"{context}: {entry!r} is not a slide number")
        if isinstance(entry, int):
            ranges.append((entry, entry))
            continue
        if isinstance(entry, str):
            match = re.fullmatch(r"\s*(\d+)\s*-\s*(\d+)\s*", entry)
            if match:
                start, end = int(match.group(1)), int(match.group(2))
                if start > end:
                    raise ValueError(f"{context}: range \"{entry}\" runs backwards")
                ranges.append((start, end))
                continue
            if entry.strip().isdigit():
                number = int(entry.strip())
                ranges.append((number, number))
                continue
        raise ValueError(f"{context}: {entry!r} is not a slide number or \"N-M\" range")
    return ranges


class SlideMap:
    """Deck-specific knowledge, loaded from JSON rather than compiled in.

    Without a map every slide lands in one output folder. That is the right
    default for a first pass over an unfamiliar deck: you cannot write the
    slide-to-page table until you have looked at what came out.
    """

    def __init__(self, data: dict[str, Any] | None, default_page: str) -> None:
        data = data or {}
        self.deck = data.get("deck")
        self.default_page = default_page
        self.pages: list[tuple[str, list[tuple[int, int]]]] = []
        self.roles: list[tuple[str, list[tuple[int, int]], list[str]]] = []

        for index, entry in enumerate(data.get("pages", [])):
            page = entry.get("page")
            if not page:
                raise ValueError(f"pages[{index}]: \"page\" is required")
            self.pages.append((page, parse_slide_selector(entry.get("slides"), f"pages[{index}]")))

        for index, entry in enumerate(data.get("roles", [])):
            role = entry.get("role")
            if not role:
                raise ValueError(f"roles[{index}]: \"role\" is required")
            match = entry.get("match", [])
            if not isinstance(match, list) or any(not isinstance(term, str) for term in match):
                raise ValueError(f"roles[{index}]: \"match\" must be a list of strings")
            slides = parse_slide_selector(entry.get("slides"), f"roles[{index}]")
            if not slides and not match:
                raise ValueError(f"roles[{index}]: needs \"slides\", \"match\", or both")
            self.roles.append((role, slides, [term.lower() for term in match]))

    @staticmethod
    def load(path: Path | None, default_page: str) -> "SlideMap":
        if path is None:
            return SlideMap(None, default_page)
        try:
            data = json.loads(path.read_text(encoding="utf-8"))
        except json.JSONDecodeError as exc:
            raise SystemExit(f"{path}: not valid JSON ({exc})") from exc
        try:
            return SlideMap(data, default_page)
        except ValueError as exc:
            raise SystemExit(f"{path}: {exc}") from exc

    def page_for_slide(self, slide_number: int) -> str:
        for page, ranges in self.pages:
            if any(start <= slide_number <= end for start, end in ranges):
                return page
        return self.default_page

    def likely_role(self, slide_number: int, text: str) -> str:
        """First matching rule wins, so order in the file is the priority order."""
        lowered = text.lower()
        for role, ranges, match in self.roles:
            if any(start <= slide_number <= end for start, end in ranges):
                return role
            if match and any(term in lowered for term in match):
                return role
        return self.page_for_slide(slide_number)


def collect_shapes(shapes: Any, parent: str | None = None) -> list[tuple[Any, str | None]]:
    collected: list[tuple[Any, str | None]] = []
    for shape in shapes:
        collected.append((shape, parent))
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            collected.extend(collect_shapes(shape.shapes, getattr(shape, "name", None)))
    return collected


def text_from_shape(shape: Any) -> str:
    if not getattr(shape, "has_text_frame", False):
        return ""
    parts: list[str] = []
    for paragraph in shape.text_frame.paragraphs:
        text = "".join(run.text for run in paragraph.runs).strip()
        if not text:
            text = paragraph.text.strip()
        if text:
            parts.append(text)
    return "\n".join(parts).strip()


def image_size(path: Path) -> dict[str, int] | None:
    try:
        with Image.open(path) as image:
            return {"width_px": image.width, "height_px": image.height}
    except Exception:
        return None


def candidate_sensitive_values(text: str) -> list[str]:
    patterns = [
        r"[A-Za-z0-9._%+-]+@[A-Za-z0-9.-]+\.[A-Za-z]{2,}",
        r"https?://[^\s)>\"]+",
        r"\b[0-9a-fA-F]{8}-[0-9a-fA-F]{4}-[0-9a-fA-F]{4}-[0-9a-fA-F]{4}-[0-9a-fA-F]{12}\b",
        r"\b(?:tenant|org|folder|client|secret|token|key)[\w -]{0,24}[:=]\s*[^\s,;]+",
    ]
    values: list[str] = []
    for pattern in patterns:
        values.extend(re.findall(pattern, text, flags=re.IGNORECASE))
    return sorted(set(values))


def callout_markers(text_blocks: list[dict], picture_blocks: list[dict]) -> list[dict]:
    """Numbered ovals in a lab deck are separate shapes with text ("1", "2", ...).

    They are lost the moment a screenshot is extracted on its own, so record each
    one as a percentage offset inside the picture it overlaps. Those percentages
    drop straight into the kit's `.workshop-click-marker` overlay in MDX.
    """
    markers = []
    for block in text_blocks:
        label = block["text"].strip()
        if not re.fullmatch(r"[0-9]{1,2}[a-z]?", label):
            continue
        rect = block["coordinates"]
        if None in (rect["x_in"], rect["y_in"], rect["w_in"], rect["h_in"]):
            continue
        cx = rect["x_in"] + rect["w_in"] / 2
        cy = rect["y_in"] + rect["h_in"] / 2
        best = None
        for picture in picture_blocks:
            pr = picture["coordinates"]
            if None in (pr["x_in"], pr["y_in"], pr["w_in"], pr["h_in"]):
                continue
            if not (pr["x_in"] <= cx <= pr["x_in"] + pr["w_in"]):
                continue
            if not (pr["y_in"] <= cy <= pr["y_in"] + pr["h_in"]):
                continue
            best = picture
            break
        markers.append(
            {
                "label": label,
                "shape_id": block["shape_id"],
                "slide_center_in": {"x": round(cx, 3), "y": round(cy, 3)},
                "on_picture": best["filename"] if best else None,
                "left_pct": (
                    round((cx - best["coordinates"]["x_in"]) / best["coordinates"]["w_in"] * 100, 1)
                    if best
                    else None
                ),
                "top_pct": (
                    round((cy - best["coordinates"]["y_in"]) / best["coordinates"]["h_in"] * 100, 1)
                    if best
                    else None
                ),
            }
        )
    return markers


def extract(
    pptx_path: Path,
    output_root: Path,
    manifest_path: Path,
    slide_map: SlideMap,
) -> dict[str, Any]:
    prs = Presentation(str(pptx_path))
    output_root.mkdir(parents=True, exist_ok=True)

    manifest: dict[str, Any] = {
        "source_pptx": str(pptx_path),
        "slide_map_deck": slide_map.deck,
        "slide_width_in": inches(prs.slide_width),
        "slide_height_in": inches(prs.slide_height),
        "slides": [],
        "sensitive_candidates": [],
        "notes": [
            "Extracted with python-pptx from slide text frames and embedded picture shapes.",
            "Coordinates are in inches from the top-left of the slide canvas.",
            "Picture files are saved from embedded image blobs; no whole-slide renders are used by this extractor.",
            "callout_markers pairs each numbered oval with the picture it overlaps, as percentage offsets for .workshop-click-marker.",
        ],
    }

    sensitive_seen: set[tuple[int, str]] = set()

    for slide_index, slide in enumerate(prs.slides, start=1):
        flat_shapes = collect_shapes(slide.shapes)
        text_blocks = []
        picture_blocks = []
        shape_notes = []

        for shape, parent in flat_shapes:
            text = text_from_shape(shape)
            if text:
                text_blocks.append(
                    {
                        "shape_id": shape.shape_id,
                        "name": shape.name,
                        "parent_group": parent,
                        "text": text,
                        "coordinates": shape_rect(shape),
                    }
                )
                for value in candidate_sensitive_values(text):
                    key = (slide_index, value)
                    if key not in sensitive_seen:
                        sensitive_seen.add(key)
                        manifest["sensitive_candidates"].append(
                            {"slide": slide_index, "value": value, "source": "slide_text"}
                        )

            # `hasattr(shape, "image")` is the whole point: an image dropped into
            # a content placeholder is a PlaceholderPicture, whose shape_type is
            # PLACEHOLDER, not PICTURE. Filtering on shape_type alone loses them.
            is_picture = shape.shape_type == MSO_SHAPE_TYPE.PICTURE or hasattr(shape, "image")
            if is_picture:
                try:
                    image = shape.image
                except (ValueError, AttributeError) as exc:
                    shape_notes.append(
                        {
                            "shape_id": shape.shape_id,
                            "name": shape.name,
                            "shape_type": str(shape.shape_type),
                            "parent_group": parent,
                            "coordinates": shape_rect(shape),
                            "note": f"Picture-shaped object has no embedded image blob: {exc}",
                        }
                    )
                    continue
                ext = image.ext or "bin"
                digest = hashlib.sha1(image.blob).hexdigest()[:10]
                page = slide_map.page_for_slide(slide_index)
                filename = f"pptx-slide-{slide_index:02d}-picture-{shape.shape_id}-{digest}.{ext}"
                out_dir = output_root / page
                out_dir.mkdir(parents=True, exist_ok=True)
                out_path = out_dir / filename
                out_path.write_bytes(image.blob)
                picture_info = {
                    "shape_id": shape.shape_id,
                    "name": shape.name,
                    "parent_group": parent,
                    "filename": str(out_path.relative_to(output_root.parent)),
                    "extension": ext,
                    "blob_sha1": hashlib.sha1(image.blob).hexdigest(),
                    "coordinates": shape_rect(shape),
                    "display_size_in": {
                        "width": inches(shape.width),
                        "height": inches(shape.height),
                    },
                    "intrinsic_size_px": image_size(out_path),
                }
                picture_blocks.append(picture_info)

            if not is_picture and not text:
                name = getattr(shape, "name", "")
                if any(token in name.lower() for token in ["callout", "arrow", "rectangle", "line", "freeform", "oval"]):
                    shape_notes.append(
                        {
                            "shape_id": shape.shape_id,
                            "name": name,
                            "shape_type": str(shape.shape_type),
                            "parent_group": parent,
                            "coordinates": shape_rect(shape),
                            "note": "Non-picture shape; may be a separate callout, box, arrow, or UI highlight.",
                        }
                    )

        text_blocks.sort(key=lambda item: (item["coordinates"]["y_in"] or 0, item["coordinates"]["x_in"] or 0))
        picture_blocks.sort(key=lambda item: (item["coordinates"]["y_in"] or 0, item["coordinates"]["x_in"] or 0))

        title = text_blocks[0]["text"] if text_blocks else ""
        subtitle = text_blocks[1]["text"] if len(text_blocks) > 1 else ""
        all_text = "\n".join(block["text"] for block in text_blocks)

        manifest["slides"].append(
            {
                "slide_number": slide_index,
                "page_assignment": slide_map.page_for_slide(slide_index),
                "likely_exercise_role": slide_map.likely_role(slide_index, all_text),
                "title": title,
                "subtitle": subtitle,
                "ordered_text_blocks": text_blocks,
                "embedded_pictures": picture_blocks,
                "separate_shape_notes": shape_notes,
                "callout_markers": callout_markers(text_blocks, picture_blocks),
            }
        )

    manifest_path.parent.mkdir(parents=True, exist_ok=True)
    manifest_path.write_text(json.dumps(manifest, indent=2), encoding="utf-8")
    return manifest


def main() -> None:
    parser = argparse.ArgumentParser(
        prog="workshop-kit extract-pptx",
        description="Extract screenshots, text and callout geometry from a facilitator deck.",
    )
    parser.add_argument("--pptx", required=True, type=Path)
    parser.add_argument("--output-root", required=True, type=Path)
    parser.add_argument("--manifest", required=True, type=Path)
    parser.add_argument(
        "--slide-map",
        type=Path,
        help="JSON mapping slides to docs pages and roles. Without it every slide "
             "goes to one folder, which is the right first pass over a new deck.",
    )
    parser.add_argument(
        "--default-page",
        default=UNASSIGNED,
        help=f"Folder for slides no page rule claims (default: {UNASSIGNED}).",
    )
    args = parser.parse_args()

    if not args.pptx.is_file():
        raise SystemExit(f"{args.pptx}: no such file")

    slide_map = SlideMap.load(args.slide_map, args.default_page)
    manifest = extract(args.pptx, args.output_root, args.manifest, slide_map)

    picture_count = sum(len(slide["embedded_pictures"]) for slide in manifest["slides"])
    text_count = sum(len(slide["ordered_text_blocks"]) for slide in manifest["slides"])
    marker_count = sum(len(slide["callout_markers"]) for slide in manifest["slides"])
    print(f"Extracted {len(manifest['slides'])} slides, {text_count} text blocks, "
          f"{picture_count} picture shapes, {marker_count} callout markers")

    # A deck run without a map, or with a map written against a different
    # revision of the deck, dumps everything into one folder. Say so, because
    # the run otherwise looks like a success.
    unassigned = [s["slide_number"] for s in manifest["slides"] if s["page_assignment"] == args.default_page]
    if unassigned:
        preview = ", ".join(str(number) for number in unassigned[:12])
        more = f" (+{len(unassigned) - 12} more)" if len(unassigned) > 12 else ""
        print(f"{len(unassigned)} slides landed in {args.default_page}/: {preview}{more}")

    print(f"Manifest: {args.manifest}")
    print(f"Image root: {args.output_root}")


if __name__ == "__main__":
    main()
